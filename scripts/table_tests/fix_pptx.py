"""Fix pptx fixture: remove hMerge fillers after gridSpan merge."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

src = Path(r'D:\one_agent\scripts\table_tests\sample.pptx')
prs = Presentation(str(src))
shape = prs.slides[0].shapes[1]  # the table shape
tbl_elem = tbl._tbl if False else shape.table._tbl
tr_last = list(tbl_elem.iter(qn('a:tr')))[-1]
tcs = tr_last.findall(qn('a:tc'))
print('Before:', len(tcs))
# tcs[0] = gridSpan=3 anchor; tcs[1], tcs[2] = hMerge fillers; tcs[3] = '240'
tr_last.remove(tcs[1])
tr_last.remove(tcs[2])
print('After:', len(tr_last.findall(qn('a:tc'))))
prs.save(str(src))
print('Saved.')
