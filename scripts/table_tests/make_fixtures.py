"""Regenerate fixtures with correct merge semantics."""
from pathlib import Path
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side
from docx import Document
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn as pqn
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

OUT = Path(__file__).resolve().parent

# ---------- xlsx ----------
wb = Workbook()
ws = wb.active
ws.title = "Sales"
ws["A1"] = "Quarterly Sales"
ws.merge_cells("A1:D1")
ws["A1"].font = Font(bold=True, size=14)
ws["A1"].alignment = Alignment(horizontal="center")
ws.append(["Region", "Q1", "Q2", "Q3"])
ws.append(["North", 120, 135, 150])
ws.append(["South", 98, 110, 90])
ws.append(["Subtotal", 218, 245, 240])
ws.append([])
ws.append(["Note: figures in USD thousands", "", "", ""])
ws.merge_cells("A7:D7")
thin = Side(style="thin", color="888888")
for row in ws["A2:D5"]:
    for cell in row:
        cell.border = Border(left=thin, right=thin, top=thin, bottom=thin)
wb.save(OUT / "sample.xlsx")

# ---------- docx ----------
doc = Document()
doc.add_heading("Quarterly Sales", level=1)
table = doc.add_table(rows=5, cols=4)
table.style = "Table Grid"
table.rows[0].cells[0].text = "Region"
table.rows[0].cells[1].text = "Q1"
table.rows[0].cells[2].text = "Q2"
table.rows[0].cells[3].text = "Q3"
table.rows[1].cells[0].text = "North"
table.rows[1].cells[1].text = "120"
table.rows[1].cells[2].text = "135"
table.rows[1].cells[3].text = "150"
table.rows[2].cells[0].text = "South"
table.rows[2].cells[1].text = "98"
table.rows[2].cells[2].text = "110"
table.rows[2].cells[3].text = "90"
r3 = table.rows[3]
r3.cells[0].text = "Subtotal (N+S)"
r3.cells[3].text = "370"
def add_gridspan(tc, n):
    tcPr = tc.get_or_add_tcPr()
    gs = OxmlElement("w:gridSpan"); gs.set(qn("w:val"), str(n)); tcPr.append(gs)
add_gridspan(r3.cells[0]._tc, 3)
# Remove the consumed filler tcs (positions 1 and 2 in the row).
tr3 = r3._tr
tcs = tr3.findall(qn("w:tc"))
tr3.remove(tcs[1])
tr3.remove(tcs[2])
r4 = table.rows[4]
r4.cells[0].text = "Total"
r4.cells[1].text = "N/A"
extra = table.add_row().cells
for c in extra: c.text = ""
def set_vmerge(tc, val):
    tcPr = tc.get_or_add_tcPr()
    vm = OxmlElement("w:vMerge")
    if val: vm.set(qn("w:val"), val)
    tcPr.append(vm)
set_vmerge(r4.cells[1]._tc, "restart")
set_vmerge(extra[1]._tc, None)
doc.add_paragraph("Notes: figures in USD thousands.")
doc.save(OUT / "sample.docx")

# ---------- pptx ----------
prs = Presentation()
prs.slide_width = Inches(10); prs.slide_height = Inches(6)
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Quarterly Sales"
shape = slide.shapes.add_table(4, 4, Inches(1), Inches(1.5), Inches(7), Inches(3))
tbl = shape.table
tbl.cell(0,0).text="Region"; tbl.cell(0,1).text="Q1"; tbl.cell(0,2).text="Q2"; tbl.cell(0,3).text="Q3"
tbl.cell(1,0).text="North"; tbl.cell(1,1).text="120"; tbl.cell(1,2).text="135"; tbl.cell(1,3).text="150"
tbl.cell(2,0).text="South"; tbl.cell(2,1).text="98"; tbl.cell(2,2).text="110"; tbl.cell(2,3).text="90"
tbl.cell(3,1).text="218"; tbl.cell(3,2).text="245"; tbl.cell(3,3).text="240"
tbl.cell(3,0).merge(tbl.cell(3,2))
tbl.cell(3,0).text = "Subtotal: all regions"
# python-pptx's merge() leaves hMerge fillers; strip them to keep the grid at 4 cols.
tbl_elem = tbl._tbl
last_tr = list(tbl_elem.iter(pqn("a:tr")))[-1]
tcs = last_tr.findall(pqn("a:tc"))
last_tr.remove(tcs[1])
last_tr.remove(tcs[2])
prs.save(OUT / "sample.pptx")

# ---------- pdf ----------
pdf_path = OUT / "sample.pdf"
doc_pdf = SimpleDocTemplate(str(pdf_path), pagesize=letter)
styles = getSampleStyleSheet()
story = []
story.append(Paragraph("Quarterly Sales Report", styles["Title"]))
story.append(Spacer(1, 12))
data1 = [["Region","Q1","Q2","Q3"],["North","120","135","150"],["South","98","110","90"],["Total","218","245","240"]]
t1 = Table(data1, hAlign="LEFT")
t1.setStyle(TableStyle([
    ("GRID", (0,0), (-1,-1), 0.5, colors.black),
    ("BACKGROUND", (0,0), (-1,0), colors.lightgrey),
    ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
]))
story.append(t1)
story.append(Spacer(1, 24))
story.append(Paragraph("This table uses visible grid lines.", styles["BodyText"]))
story.append(PageBreak())
story.append(Paragraph("Whitespace-Aligned Table", styles["Heading2"]))
story.append(Spacer(1, 12))
mono = ParagraphStyle("mono", parent=styles["BodyText"], fontName="Courier", fontSize=10, leading=12)
data2 = [["Item","Jan","Feb","Mar"],["Widgets","120","130","145"],["Gizmos","85","92","78"],["Sprockets","44","50","55"]]
data2_mono = [[Paragraph(c, mono) for c in row] for row in data2]
t2 = Table(data2_mono, hAlign="LEFT", colWidths=[110,60,60,60])
story.append(t2)
story.append(Spacer(1, 12))
story.append(Paragraph("End of report.", styles["BodyText"]))
doc_pdf.build(story)

for p in [OUT / "sample.xlsx", OUT / "sample.docx", OUT / "sample.pptx", OUT / "sample.pdf"]:
    print(f"  {p}  ({p.stat().st_size} bytes)")
