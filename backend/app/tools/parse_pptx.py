"""Parse .pptx -> Markdown text via python-pptx, preserving merged cells.

PPTX tables use DrawingML:
  - gridSpan="N" on the anchor -> spans N columns; the (N-1) following cells
    in the same row carry hMerge="true" and are blank.
  - rowSpan="N" on the anchor -> spans N rows; the (N-1) following cells in
    the same column carry vMerge="true" and are blank.
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.oxml.ns import qn


def _md_escape(s: str) -> str:
    return s.replace("|", "\\|").replace("\n", " ").replace("\r", " ")


def _attr(elem, name: str) -> str | None:
    return elem.get(name)


def _int_attr(el, name: str) -> int | None:
    v = _attr(el, name)
    if v is None:
        return None
    try:
        return int(v)
    except (TypeError, ValueError):
        return None


def _bool_attr(el, name: str) -> bool:
    v = _attr(el, name)
    return v in ("true", "1")


def _tc_grid_span(tc) -> int:
    return _int_attr(tc, "gridSpan") or 1


def _tc_row_span(tc) -> int:
    return _int_attr(tc, "rowSpan") or 1


def _tc_text(tc) -> str:
    parts = []
    for t in tc.iter(qn("a:t")):
        parts.append(t.text or "")
    return "".join(parts).strip()


def _is_hmerge(tc) -> bool:
    return _bool_attr(tc, "hMerge")


def _is_vmerge(tc) -> bool:
    return _bool_attr(tc, "vMerge")


def _grid_to_markdown(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    n_cols = max((len(r) for r in rows), default=0)
    rows = [r + [""] * (n_cols - len(r)) for r in rows]
    while rows and not any(c.strip() for c in rows[-1]):
        rows.pop()
    if not rows:
        return ""
    widths = [0] * n_cols
    for r in rows:
        for i, c in enumerate(r):
            widths[i] = max(widths[i], len(_md_escape(c)))

    def fmt_row(r):
        return "| " + " | ".join(_md_escape(c).ljust(widths[i]) for i, c in enumerate(r)) + " |"

    header = rows[0]
    sep = "| " + " | ".join("-" * widths[i] for i in range(n_cols)) + " |"
    body = "\n".join(fmt_row(r) for r in rows[1:])
    out = fmt_row(header) + "\n" + sep
    if body:
        out += "\n" + body
    return out


def _shape_table_to_markdown(shape) -> str:
    """Convert a pptx shape with a table into Markdown."""
    tbl = shape.table._tbl  # underlying a:tbl element
    # Track vertical-merge anchors per column so continuation rows can show
    # the anchor's text (matches how Word/pptx renders the merged cell).
    vmerge_anchor: dict[int, str] = {}
    rows_out: list[list[str]] = []

    for tr in tbl.iter(qn("a:tr")):
        cells: list[str] = []
        col_idx = 0
        for tc in tr.findall(qn("a:tc")):
            grid_span = _tc_grid_span(tc)
            row_span = _tc_row_span(tc)
            hmerge = _is_hmerge(tc)
            vmerge = _is_vmerge(tc)

            if hmerge:
                # Continuation of a horizontal merge -> blank.
                cells.append("")
                for _ in range(grid_span - 1):
                    cells.append("")
                col_idx += grid_span
                continue

            if vmerge:
                # Continuation of a vertical merge -> emit anchor's text.
                anchor = vmerge_anchor.get(col_idx, "")
                cells.append(anchor)
                for _ in range(grid_span - 1):
                    cells.append("")
                col_idx += grid_span
                continue

            text = _tc_text(tc)
            cells.append(text)
            for _ in range(grid_span - 1):
                cells.append("")
            if row_span > 1:
                vmerge_anchor[col_idx] = text
            col_idx += grid_span
        rows_out.append(cells)

    return _grid_to_markdown(rows_out)


def parse_pptx(path: str) -> dict:
    if not os.path.isfile(path):
        raise FileNotFoundError(path)
    pres = Presentation(path)
    slides = []
    for i, slide in enumerate(pres.slides, 1):
        parts = []
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(slide.shapes.title.text.strip())
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        parts.append(t)
            elif shape.has_table:
                md = _shape_table_to_markdown(shape)
                if md:
                    parts.append(md)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[notes] {notes}")
        if parts:
            slides.append(f"[Slide {i}]\n" + "\n".join(parts))
    content = "\n\n".join(slides).strip()
    title = os.path.splitext(os.path.basename(path))[0] or "Untitled Slides"
    if not content:
        content = "(empty slides)"
    return {"title": title[:200], "content": content}
